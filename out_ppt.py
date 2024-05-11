from pptx import Presentation
from pptx.slide import Slide
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE
from enum import Enum
import pandas as pd
from load_csv import get_institution_logo, get_institution_team, get_institution_debater, get_institution_adjudicator

class BlankSlideSettings(Enum):
    NONE = "None"
    ALL = "All"
    CHANGE_RANK = "Change of rankings"

def list_layouts(pres: Presentation) -> None:
    for index, layout in enumerate(pres.slide_layouts):
        print(f"[{index}] {layout.name}")

def add_result_slide(pres: Presentation, layout_index: int, title: str, metrics: str, name: str, logos: list) -> Slide:
    pictindex = 0
    slide_layout = pres.slide_layouts[layout_index]
    slide = pres.slides.add_slide(slide_layout)
    for shape_layout in slide_layout.placeholders:
        try:
            shape_slide = slide.placeholders[shape_layout.placeholder_format.idx]
            if shape_layout.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                if (len(logos) > pictindex):
                    pic = shape_slide.insert_picture(logos[pictindex])
                    pict_fit(pic)
                    pictindex += 1
                    
            else:
                shape_slide.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                shape_slide.text_frame.text = shape_layout.text.format(title = title, metrics = metrics, name = name)
                #shape_slide.text_frame.fit_text()
        except Exception:
            continue

def pict_fit(picture):
    crop_v = picture.crop_top + picture.crop_bottom
    crop_h = picture.crop_left + picture.crop_right
    picture.crop_left = -crop_v / (1 - crop_v) / 2
    picture.crop_right = -crop_v / (1 - crop_v) / 2
    picture.crop_top = -crop_h / (1 - crop_h) / 2
    picture.crop_bottom = -crop_h / (1 - crop_h) / 2

def create_slides_teams(pres, dict_slides: dict, df_standings: pd.DataFrame, df_teams: pd.DataFrame, dict_logo: dict, title_text: list, metrics_show: list, metrics_hide: list, dict_display: dict, danger_prevention: BlankSlideSettings = BlankSlideSettings.CHANGE_RANK):
    metrics_all = [*metrics_show, *metrics_hide]
    value_counts = df_standings["rank"].value_counts()
    subrank = 1
    prev_rank = 0
    for index, row in df_standings.iterrows():
        paths_institution = get_institution_logo(dict_logo, get_institution_team(df_teams, row["team"]))
        # Set title
        str_title = title_text[0 if row["rank"] > 0 else 1].format(ordinal(abs(int(row["rank"])))).replace("1st Best", "Best")
        if row["rank"] == prev_rank:
            subrank += 1
        else:
            subrank = 1
        prev_rank = row["rank"]
        if value_counts[row['rank']] >= 2:
            str_title += f" ({subrank}/{value_counts[row['rank']]})"
        # Set metrics
        str_metric = ", ".join([dict_display[metric].format(row[metric.value]) for metric in metrics_all if pd.notna(row[metric.value])])
        # Danger prevention slides
        if danger_prevention == BlankSlideSettings.ALL or (danger_prevention == BlankSlideSettings.CHANGE_RANK and subrank == 1):
            add_result_slide(pres, dict_slides[0], str_title, str_metric, "", [])
        # Real Slide
        add_result_slide(pres, dict_slides[len(paths_institution)], str_title, str_metric, row["team"], paths_institution)

def create_slides_debaters(pres, dict_slides: dict, df_standings: pd.DataFrame, df_teams: pd.DataFrame, dict_logo: pd.DataFrame, title_text: list, metrics_show: list, metrics_hide: list, dict_metrics: dict, danger_prevention: BlankSlideSettings = BlankSlideSettings.CHANGE_RANK):
    metrics_all = [*metrics_show, *metrics_hide]
    value_counts = df_standings["rank"].value_counts()
    subrank = 1
    prev_rank = 0
    for index, row in df_standings.iterrows():
        paths_institution = get_institution_logo(dict_logo, get_institution_debater(df_teams, row["name"]))
        # Set title
        str_title = title_text.format(ordinal(abs(int(row["rank"])))).replace("1st Best", "Best")
        if row["rank"] == prev_rank:
            subrank += 1
        else:
            subrank = 1
        prev_rank = row["rank"]
        if value_counts[row['rank']] >= 2:
            str_title += f" ({subrank}/{value_counts[row['rank']]})"
        # Set metrics
        str_metric = ", ".join([dict_metrics[metric].format(row[metric.value]) for metric in metrics_all if pd.notna(row[metric.value])])
        # Danger prevention slides
        if danger_prevention == BlankSlideSettings.ALL or (danger_prevention == BlankSlideSettings.CHANGE_RANK and subrank == 1):
            add_result_slide(pres, dict_slides[0], str_title, str_metric, "", [])
        # Real Slide
        add_result_slide(pres, dict_slides[len(paths_institution)], str_title, str_metric, row["name"], paths_institution)

def create_slides_adjudicators(pres, dict_slides: dict, df_standings: pd.DataFrame, df_adjs: pd.DataFrame, dict_logo: dict, title_text: list, display_score: str, danger_prevention: BlankSlideSettings = BlankSlideSettings.CHANGE_RANK):
    value_counts = df_standings["rank"].value_counts()
    subrank = 1
    prev_rank = 0
    for index, row in df_standings.iterrows():
        paths_institution = get_institution_logo(dict_logo, get_institution_adjudicator(df_adjs, row["name"]))
        # Set title
        str_title = title_text.format(ordinal(abs(int(row["rank"])))).replace("1st Best", "Best")
        if row["rank"] == prev_rank:
            subrank += 1
        else:
            subrank = 1
        prev_rank = row["rank"]
        if value_counts[row['rank']] >= 2:
            str_title += f" ({subrank}/{value_counts[row['rank']]})"
        # Set metrics
        str_metric = display_score.format(row["score"])
        # Danger prevention slides
        if danger_prevention == BlankSlideSettings.ALL or (danger_prevention == BlankSlideSettings.CHANGE_RANK and subrank == 1):
            add_result_slide(pres, dict_slides[0], str_title, str_metric, "", [])
        # Real Slide
        add_result_slide(pres, dict_slides[len(paths_institution)], str_title, str_metric, row["name"], paths_institution)

def ordinal(n: int):
    if 11 <= (n % 100) <= 13:
        suffix = 'th'
    else:
        suffix = ['th', 'st', 'nd', 'rd', 'th'][min(n % 10, 4)]
    return str(n) + suffix
